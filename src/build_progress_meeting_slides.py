from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parents[1]
OUTPUT = BASE / "reports" / "slides" / "progress_meeting_course_clustering.pptx"
REPORT_FIGURES = BASE / "results" / "figures" / "keep_for_report"
DENDROGRAM = REPORT_FIGURES / "hierarchical_dendrogram.png"
HEATMAP = REPORT_FIGURES / "course_similarity_heatmap.png"
ASSIGNMENTS = BASE / "results" / "tables" / "hierarchical_cluster_assignments.csv"
SUMMARY = BASE / "results" / "tables" / "cluster_summary.csv"
SIMILARITY = BASE / "results" / "tables" / "course_similarity_matrix.csv"

NAVY = RGBColor(31, 43, 61)
TEAL = RGBColor(32, 135, 130)
GREEN = RGBColor(83, 145, 96)
AMBER = RGBColor(214, 151, 68)
BURGUNDY = RGBColor(145, 67, 72)
INK = RGBColor(39, 43, 50)
MUTED = RGBColor(97, 106, 116)
LIGHT = RGBColor(244, 246, 248)
SOFT_TEAL = RGBColor(224, 242, 238)
SOFT_AMBER = RGBColor(249, 239, 220)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(214, 219, 225)


def add_footer(slide, number: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = f"NOVA50301 AI Toolkit · Progress Meeting · {number}/5"
    p.font.name = "Aptos"
    p.font.size = Pt(8)
    p.font.color.rgb = MUTED


def add_title(slide, title: str, subtitle: str, number: int) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(9.8), Inches(0.45))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.75), Inches(9.8), Inches(0.28))
    p = sub_box.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.name = "Aptos"
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(11.7), Inches(0.32), Inches(1.05), Inches(0.38))
    tag.fill.solid()
    tag.fill.fore_color.rgb = NAVY
    tag.line.fill.background()
    p = tag.text_frame.paragraphs[0]
    p.text = f"{number} / 5"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_section_label(slide, text: str, x: float, y: float, w: float, color: RGBColor = TEAL) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_card(slide, x: float, y: float, w: float, h: float, title: str, body: str, fill: RGBColor = LIGHT) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(4)
    p = frame.add_paragraph()
    p.text = body
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = INK


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, size: int = 13) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(5)


def add_pipeline(slide) -> None:
    steps = [
        ("1", "학과 선정", "25개 목적 표본"),
        ("2", "데이터 구성", "학과×과목 binary matrix"),
        ("3", "유사도 계산", "Cosine similarity"),
        ("4", "군집화", "Hierarchical clustering"),
        ("5", "해석", "상담 지원용 예비 근거"),
    ]
    x = 0.65
    colors = [TEAL, GREEN, AMBER, BURGUNDY, NAVY]
    for idx, (num, title, body) in enumerate(steps):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.15), Inches(2.15), Inches(0.88))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[idx]
        shape.line.fill.background()
        frame = shape.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.text = num
        p.font.name = "Aptos"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p = frame.add_paragraph()
        p.text = title
        p.font.name = "Aptos"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p = frame.add_paragraph()
        p.text = body
        p.font.name = "Aptos"
        p.font.size = Pt(8.5)
        p.font.color.rgb = WHITE
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x + 2.22), Inches(5.43), Inches(0.34), Inches(0.28))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LINE
            arrow.line.fill.background()
        x += 2.52


def add_matrix_preview(slide) -> None:
    rows = [
        ["department", "미적분Ⅰ", "물리학", "화학", "생명과학", "경제"],
        ["기계공학과", "1", "1", "1", "1", "0"],
        ["자동차공학과", "1", "1", "1", "0", "0"],
        ["경제학과", "1", "0", "0", "0", "1"],
    ]
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.72), Inches(3.75), Inches(6.25), Inches(1.52)).table
    widths = [1.75, 0.9, 0.9, 0.9, 0.95, 0.85]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Aptos"
            p.font.size = Pt(8.3)
            p.font.color.rgb = WHITE if r == 0 else INK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                p.font.bold = True
            elif text == "1":
                cell.fill.solid()
                cell.fill.fore_color.rgb = SOFT_TEAL


def add_image(slide, path: Path, x: float, y: float, w: float, h: float, caption: str) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT
        shape.line.color.rgb = LINE
        p = shape.text_frame.paragraphs[0]
        p.text = f"Missing figure: {path.name}"
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = MUTED
    box = slide.shapes.add_textbox(Inches(x), Inches(y + h + 0.05), Inches(w), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = caption
    p.font.name = "Aptos"
    p.font.size = Pt(8.5)
    p.font.color.rgb = MUTED


def cluster_counts() -> str:
    if not ASSIGNMENTS.exists():
        return "cluster assignments not generated"
    df = pd.read_csv(ASSIGNMENTS, encoding="utf-8-sig")
    counts = df.groupby("cluster_id")["department_name"].apply(list)
    parts = []
    for cluster_id, names in counts.items():
        label = ", ".join(names[:3])
        suffix = "..." if len(names) > 3 else ""
        parts.append(f"C{cluster_id}: {len(names)}개 ({label}{suffix})")
    return "\n".join(parts)


def top_pairs(limit: int = 4) -> list[str]:
    if not SIMILARITY.exists():
        return ["similarity matrix not generated"]
    sim = pd.read_csv(SIMILARITY, encoding="utf-8-sig", index_col=0)
    ids = [column for column in sim.columns if column != "department_name"]
    names = sim["department_name"].to_dict()
    rows = []
    for i, left in enumerate(ids):
        for right in ids[i + 1 :]:
            rows.append((names[left], names[right], float(sim.loc[left, right])))
    rows.sort(key=lambda row: row[2], reverse=True)
    return [f"{left}–{right}: {value:.3f}" for left, right, value in rows[:limit]]


def make_slide(prs: Presentation, title: str, subtitle: str, number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_title(slide, title, subtitle, number)
    add_footer(slide, number)
    return slide


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = make_slide(prs, "권장과목 기반 학과 군집화", "입시 상담 지원을 위한 탐색적 분석 · Progress Update", 1)
    add_section_label(slide, "핵심 질문", 0.65, 1.35, 1.05, TEAL)
    add_card(slide, 0.65, 1.78, 5.55, 1.15, "Research Question", "권장과목 profile이 학과들을 해석 가능한 방식으로 군집화할 수 있는가?", WHITE)
    add_section_label(slide, "Scope", 6.65, 1.35, 0.85, NAVY)
    add_card(slide, 6.65, 1.78, 5.95, 1.15, "Progress-stage scope", "탐색적 군집화 연구 (학과 추천 시스템 X)\n25개 학과 목적 표본 · binary vector · hierarchical clustering", LIGHT)
    add_section_label(slide, "Pipeline", 0.65, 4.65, 0.95, AMBER)
    add_pipeline(slide)
    slide.notes_slide.notes_text_frame.text = "첫 장에서는 추천 시스템이 아니라 탐색적 군집화 연구라는 점을 먼저 고정한다."

    slide = make_slide(prs, "데이터와 표본", "25개 학과 · 학과×과목 행렬 · binary 인코딩", 2)
    add_section_label(slide, "학과 × 권장과목 행렬", 0.72, 1.28, 1.9, TEAL)
    add_bullets(slide, ["원자료: 학과 과목 선택 가이드.xlsx", "A열 학과명 · B-D열 일반선택/진로선택/융합선택", "1 = 가이드에 제시됨 · 0 = 제시되지 않음", "metadata columns는 clustering feature에서 제외"], 0.72, 1.78, 5.9, 1.55, 12)
    add_matrix_preview(slide)
    add_section_label(slide, "25개 학과 선정 근거", 7.35, 1.28, 1.75, NAVY)
    add_card(slide, 7.35, 1.78, 4.9, 0.72, "목적 표본", "통계적 대표 표본이 아니라 exploratory purposive sample", SOFT_TEAL)
    add_card(slide, 7.35, 2.68, 4.9, 0.72, "권장과목 다양성", "인문, 사회, 공학, 자연, 의약/보건, 디자인을 함께 포함", LIGHT)
    add_card(slide, 7.35, 3.58, 4.9, 0.72, "울산 산업 boundary case", "조선해양공학과와 자동차공학과를 applied engineering case로 포함", SOFT_AMBER)
    add_card(slide, 7.35, 4.48, 4.9, 0.72, "해석 가능성", "25개 규모로 heatmap과 dendrogram을 직접 해석 가능", LIGHT)
    slide.notes_slide.notes_text_frame.text = "표본은 대표성이 아니라 비교 가능성과 해석 가능성을 위해 고른 목적 표본이다."

    slide = make_slide(prs, "방법과 군집 구조", "Cosine similarity + Hierarchical clustering", 3)
    add_section_label(slide, "방법 선택 근거", 0.65, 1.22, 1.45, TEAL)
    add_card(slide, 0.65, 1.68, 5.15, 0.83, "Cosine similarity", "학과별 과목 수 차이가 있으므로 magnitude보다 course composition을 비교", LIGHT)
    add_card(slide, 0.65, 2.68, 5.15, 0.83, "Hierarchical clustering", "n=25 소표본에서 dendrogram으로 구조를 직접 확인하기 적합", LIGHT)
    add_card(slide, 0.65, 3.68, 5.15, 0.83, "Feature rule", "course columns만 사용하고 department_id/name/broad_field는 제외", SOFT_TEAL)
    add_image(slide, DENDROGRAM, 6.25, 1.35, 6.45, 4.15, "Figure 1 · Dendrogram (Average linkage, Cosine distance)")
    add_card(slide, 6.25, 5.85, 6.45, 0.55, "군집 요약", cluster_counts(), WHITE)
    slide.notes_slide.notes_text_frame.text = "dendrogram은 hard clustering보다 학과 간 가까움과 멀어짐을 설명하는 보조 도구로 사용한다."

    slide = make_slide(prs, "예비 결과", "Heatmap · 상위 유사도 pair · 해석 포인트", 4)
    add_image(slide, HEATMAP, 0.65, 1.25, 6.25, 4.75, "Figure 2 · Course-profile cosine similarity heatmap")
    add_section_label(slide, "정량 요약", 7.25, 1.25, 1.1, NAVY)
    add_card(slide, 7.25, 1.72, 4.95, 1.25, "Top similarity pairs", "\n".join(top_pairs()), SOFT_TEAL)
    add_card(slide, 7.25, 3.18, 4.95, 0.95, "핵심 발견 1", "자동차공학과는 기계공학과·전기전자공학과와 높은 유사도를 보여 applied engineering boundary case로 설명 가능", LIGHT)
    add_card(slide, 7.25, 4.3, 4.95, 0.95, "핵심 발견 2", "공학/자연/의약/정량 계열이 큰 군집으로 묶여 세부 해석에는 local similarity 확인이 필요", LIGHT)
    add_card(slide, 7.25, 5.42, 4.95, 0.72, "주의", "결과는 exploratory pattern이며 최종 추천/분류가 아님", SOFT_AMBER)
    slide.notes_slide.notes_text_frame.text = "예비 결과는 예쁘게 잘린 군집보다 boundary case와 local similarity를 보여주는 데 의미가 있다."

    slide = make_slide(prs, "다음 단계", "Final Report에서 확장할 분석과 보수적 해석", 5)
    add_section_label(slide, "WHY", 0.65, 1.25, 0.7, BURGUNDY)
    add_card(slide, 0.65, 1.7, 11.9, 0.9, "현재 결과의 의미", "Hierarchical clustering은 권장과목 profile의 예비 구조를 보여주지만, 큰 군집과 singleton cluster의 해석에는 민감도 분석과 제한점 논의가 필요하다.", WHITE)
    add_section_label(slide, "Final Report Roadmap", 0.65, 3.0, 2.0, NAVY)
    add_card(slide, 0.65, 3.48, 2.75, 1.0, "W13", "Progress 발표\n현재 결과 점검", SOFT_TEAL)
    add_card(slide, 3.65, 3.48, 2.75, 1.0, "W14", "k-means 비교\ncluster 수 민감도", LIGHT)
    add_card(slide, 6.65, 3.48, 2.75, 1.0, "W15", "결과 해석 정리\n한계와 boundary case", LIGHT)
    add_card(slide, 9.65, 3.48, 2.75, 1.0, "W16", "5-page final report\nFuture work 정리", SOFT_AMBER)
    add_section_label(slide, "Future Work", 0.65, 5.25, 1.35, TEAL)
    add_bullets(slide, ["Expert consensus, admission score feasibility, candidate generation은 future work", "현재 Progress 단계는 department-course matrix와 clustering 결과에 제한", "최종 주장은 '자동 추천'이 아니라 '상담 지원을 위한 해석 가능한 유사성 탐색'"], 0.75, 5.72, 11.4, 0.8, 12)
    slide.notes_slide.notes_text_frame.text = "마지막 장에서는 범위 확장을 막고 final report에서 할 일만 보수적으로 제시한다."

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
