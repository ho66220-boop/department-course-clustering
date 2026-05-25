from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(__file__).resolve().parents[1]
OUTPUT = BASE / "reports" / "slides" / "progress_meeting_course_clustering.pptx"
DENDROGRAM = BASE / "results" / "figures" / "hierarchical_dendrogram.png"
HEATMAP = BASE / "results" / "figures" / "course_similarity_heatmap.png"

NAVY = RGBColor(32, 45, 63)
TEAL = RGBColor(39, 142, 137)
GREEN = RGBColor(92, 151, 96)
AMBER = RGBColor(214, 151, 68)
BURGUNDY = RGBColor(145, 67, 72)
INK = RGBColor(38, 42, 48)
MUTED = RGBColor(99, 108, 118)
LIGHT = RGBColor(244, 246, 248)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(214, 219, 225)


def add_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.55))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(0.88),
        Inches(1.15),
        Inches(0.06),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def add_footer(slide, number: int) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.1), Inches(0.22))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = f"NOVA50301 AI Toolkit Progress Meeting | {number}/5"
    p.font.name = "Aptos"
    p.font.size = Pt(8)
    p.font.color.rgb = MUTED


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(17)
        p.font.color.rgb = INK
        p.space_after = Pt(7)


def add_label(slide, text: str, x: float, y: float, w: float, h: float, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_caption(slide, text: str, x: float, y: float, w: float, h: float) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def add_notes(slide, notes: str) -> None:
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes


def add_pipeline(slide) -> None:
    steps = [
        ("Department\nselection", TEAL),
        ("Course-profile\nmatrix", GREEN),
        ("Hierarchical\nclustering", AMBER),
        ("Exploratory\ninterpretation", BURGUNDY),
    ]
    x = 0.85
    for idx, (text, color) in enumerate(steps):
        add_label(slide, text, x, 5.2, 2.3, 0.68, color)
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                Inches(x + 2.42),
                Inches(5.34),
                Inches(0.48),
                Inches(0.32),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LINE
            arrow.line.fill.background()
        x += 3.05
    add_caption(slide, "Recommended figure: project pipeline diagram", 0.85, 6.05, 9.0, 0.25)


def add_selection_table(slide) -> None:
    table = slide.shapes.add_table(4, 3, Inches(0.85), Inches(4.35), Inches(11.6), Inches(1.55)).table
    widths = [3.0, 4.2, 4.4]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    headers = ["Selection Criterion", "Meaning", "Role in Project"]
    rows = [
        ["Course-profile diversity", "Include contrasting course patterns", "Improve clustering signal"],
        ["Counseling relevance", "Reflect realistic guidance contexts", "Support admission counseling use"],
        ["Interpretability", "Keep sample explainable", "Fit a short exploratory study"],
    ]
    for col, text in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = text
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
    for row_idx, row in enumerate(rows, start=1):
        for col, text in enumerate(row):
            cell = table.cell(row_idx, col)
            cell.text = text
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.color.rgb = INK
    add_caption(slide, "Recommended table: rationale for purposive department selection", 0.85, 6.05, 9.8, 0.25)


def add_matrix_table(slide) -> None:
    table = slide.shapes.add_table(4, 6, Inches(0.85), Inches(4.05), Inches(11.6), Inches(1.85)).table
    widths = [3.2, 1.55, 1.55, 1.55, 1.55, 1.55]
    for idx, width in enumerate(widths):
        table.columns[idx].width = Inches(width)
    rows = [
        ["department_name", "Mathematics", "Physics", "Chemistry", "Biology", "Economics"],
        ["Mechanical Engineering", "1.0", "1.0", "0.5", "0.0", "0.0"],
        ["Biology", "0.5", "0.0", "1.0", "1.0", "0.0"],
        ["Business Administration", "0.5", "0.0", "0.0", "0.0", "1.0"],
    ]
    for row_idx, row in enumerate(rows):
        for col, text in enumerate(row):
            cell = table.cell(row_idx, col)
            cell.text = text
            cell.text_frame.paragraphs[0].font.size = Pt(8.5)
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE if row_idx == 0 else INK
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                cell.text_frame.paragraphs[0].font.bold = True
            elif text == "1.0":
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(224, 242, 238)
            elif text == "0.5":
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(246, 237, 214)
    add_caption(slide, "Recommended table: department-course matrix preview", 0.85, 6.05, 8.5, 0.25)


def add_result_figure(slide) -> None:
    image = DENDROGRAM if DENDROGRAM.exists() else HEATMAP
    if image.exists():
        slide.shapes.add_picture(str(image), Inches(6.75), Inches(1.35), width=Inches(5.9), height=Inches(4.2))
        add_caption(slide, f"Recommended figure: {image.name}", 6.75, 5.7, 5.6, 0.32)
    else:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.75), Inches(1.35), Inches(5.9), Inches(4.2))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = LINE
        frame = box.text_frame
        frame.text = "Insert dendrogram or heatmap"
        frame.paragraphs[0].font.color.rgb = MUTED
        frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_roadmap_table(slide) -> None:
    table = slide.shapes.add_table(5, 3, Inches(0.85), Inches(4.05), Inches(11.6), Inches(1.95)).table
    headers = ["Stage", "Task", "Purpose"]
    rows = [
        ["Current", "Hierarchical clustering", "Initial exploratory grouping"],
        ["Next", "K-means comparison", "Check method stability"],
        ["Next", "Binary-vector sensitivity", "Test value-coding sensitivity"],
        ["Optional", "Expert-based comparison", "Validate interpretability"],
    ]
    for col, text in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = text
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
    for row_idx, row in enumerate(rows, start=1):
        for col, text in enumerate(row):
            cell = table.cell(row_idx, col)
            cell.text = text
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.color.rgb = INK
    add_caption(slide, "Recommended table: final-project analysis roadmap", 0.85, 6.12, 8.5, 0.25)


def make_slide(prs, number: int, title: str, bullets: list[str], notes: str) -> object:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_title(slide, title)
    add_bullets(slide, bullets, 0.85, 1.35, 5.45, 3.15)
    add_footer(slide, number)
    add_notes(slide, notes)
    return slide


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = make_slide(
        prs,
        1,
        "Project Overview",
        [
            "Exploratory clustering of academic departments using recommended high-school course profiles",
            "Designed to support admission counseling, not to replace counseling decisions",
            "The project is not a complete department recommendation system",
            "Main question: can course profiles produce interpretable department clusters?",
            "Current stage focuses on data structure and preliminary clustering results",
        ],
        "이 프로젝트의 핵심은 학생에게 어떤 학과를 추천할 것인가를 바로 결정하는 시스템이 아닙니다. 현재 목표는 학과별 추천 고등학교 과목 정보를 사용했을 때, 학과들이 의미 있게 묶이는지 확인하는 탐색적 clustering입니다. 결과는 상담을 보조하기 위한 분석 자료로 보는 것이 적절합니다.",
    )
    add_pipeline(slide)

    slide = make_slide(
        prs,
        2,
        "Why 24 Departments?",
        [
            "The 24 departments are a purposive sample, not a statistically representative sample",
            "Departments were selected to maximize course-profile diversity",
            "Counseling relevance was considered to reflect realistic admission guidance contexts",
            "Interpretability was prioritized for a short exploratory clustering study",
            "The sample size was chosen to keep the analysis manageable and explainable",
        ],
        "24개 학과는 전체 학과를 통계적으로 대표하기 위한 무작위 표본이 아닙니다. 짧은 기간 안에 clustering 가능성과 해석 가능성을 확인하기 위해, 과목 프로필 다양성, 상담 관련성, 해석 가능성을 기준으로 목적 표본을 구성했습니다.",
    )
    add_selection_table(slide)

    slide = make_slide(
        prs,
        3,
        "Data Structure",
        [
            "Main dataset: a department-course matrix",
            "Rows represent academic departments",
            "Columns represent recommended high-school courses",
            "Course values: 1.0 = core, 0.5 = related/supporting, 0.0 = not mentioned",
            "Only course columns are used as clustering features",
        ],
        "분석의 핵심 데이터는 department-course matrix입니다. 각 행은 학과이고 각 열은 추천 과목입니다. 값은 핵심 추천 과목 1.0, 관련 또는 보조 과목 0.5, 언급되지 않은 과목 0.0으로 구성했습니다. clustering에는 과목 열만 사용하고 metadata는 feature에서 제외했습니다.",
    )
    add_matrix_table(slide)

    slide = make_slide(
        prs,
        4,
        "Method and Preliminary Result",
        [
            "Hierarchical clustering was implemented as the first clustering method",
            "Clustering is based only on course-profile similarity",
            "Preliminary outputs include similarity matrix, heatmap, dendrogram, and cluster assignments",
            "Results should be interpreted as exploratory, not definitive",
            "Early clusters help check whether course-based patterns are interpretable",
        ],
        "현재 구현된 첫 번째 방법은 hierarchical clustering입니다. 학과 간 유사도는 학과명이나 계열 정보가 아니라 추천 과목 프로필만으로 계산됩니다. similarity matrix, heatmap, dendrogram, cluster assignment를 예비 결과로 볼 수 있습니다. 다만 이 결과는 최종 분류가 아니라 exploratory result로 해석해야 합니다.",
    )
    add_result_figure(slide)

    slide = make_slide(
        prs,
        5,
        "Next Steps",
        [
            "Compare hierarchical clustering with k-means clustering",
            "Conduct sensitivity analysis using binary course vectors",
            "Refine interpretation of cluster characteristics",
            "Optionally compare data-driven clusters with expert-based grouping",
            "Position expert consensus and admission feasibility as future work",
        ],
        "다음 단계에서는 hierarchical clustering 결과만으로 결론을 내리지 않고 k-means clustering과 비교해 군집 구조의 안정성을 확인할 계획입니다. 또한 binary vector를 사용한 sensitivity analysis를 수행할 수 있습니다. expert consensus와 admission feasibility는 future work로 남겨두겠습니다.",
    )
    add_roadmap_table(slide)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
